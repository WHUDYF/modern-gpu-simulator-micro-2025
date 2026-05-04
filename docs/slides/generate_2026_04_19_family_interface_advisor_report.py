from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDES_DIR = Path(__file__).resolve().parent
OUT_PATH = SLIDES_DIR / "2026-04-19-family-interface-advisor-report.pptx"
ASSET_DIR = SLIDES_DIR / "assets"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG = RGBColor(247, 248, 250)
TITLE_BG = RGBColor(22, 35, 58)
ACCENT = RGBColor(201, 76, 76)
ACCENT_2 = RGBColor(47, 128, 147)
TEXT = RGBColor(33, 37, 41)
SUBTLE = RGBColor(96, 103, 112)
BOX_BG = RGBColor(255, 255, 255)
BOX_LINE = RGBColor(215, 220, 226)
HILITE_BG = RGBColor(236, 243, 247)
WARN_BG = RGBColor(251, 240, 240)

FONT = "Microsoft YaHei"
TITLE_FONT = "Microsoft YaHei"


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_top_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.0)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TITLE_BG
    bar.line.color.rgb = TITLE_BG

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.22), Inches(10.0), Inches(0.38))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = TITLE_FONT
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.62), Inches(0.63), Inches(11.0), Inches(0.2))
        p = sub_box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = FONT
        r.font.size = Pt(11)
        r.font.color.rgb = RGBColor(215, 224, 235)


def add_footer(slide, note):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(12.0), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = note
    r.font.name = FONT
    r.font.size = Pt(10)
    r.font.color.rgb = SUBTLE


def add_image_card(slide, left, top, width, height, img_path, title, bullets, citation):
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = BOX_BG
    panel.line.color.rgb = BOX_LINE

    add_text_box(slide, left + Inches(0.2), top + Inches(0.12), width - Inches(0.4), Inches(0.3),
                 title, font_size=18, bold=True)

    pic_left = left + Inches(0.18)
    pic_top = top + Inches(0.55)
    pic_w = width * 0.52
    pic_h = height - Inches(1.0)
    slide.shapes.add_picture(str(img_path), pic_left, pic_top, width=pic_w, height=pic_h)

    body = slide.shapes.add_textbox(left + width * 0.58, top + Inches(0.55), width * 0.35, height - Inches(1.15))
    tf = body.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(6)
        p.font.name = FONT
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT

    cite = slide.shapes.add_textbox(left + Inches(0.2), top + height - Inches(0.33), width - Inches(0.4), Inches(0.18))
    p = cite.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = citation
    r.font.name = FONT
    r.font.size = Pt(10)
    r.font.color.rgb = SUBTLE


def add_text_box(slide, left, top, width, height, text, font_size=20, bold=False,
                 color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_bullet_panel(slide, left, top, width, height, heading, bullets, fill_color=BOX_BG):
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = fill_color
    panel.line.color.rgb = BOX_LINE

    add_text_box(slide, left + Inches(0.2), top + Inches(0.12), width - Inches(0.4), Inches(0.3),
                 heading, font_size=18, bold=True)

    body = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.5), width - Inches(0.35), height - Inches(0.6))
    tf = body.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(5)
        p.font.name = FONT
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT


def add_banner(slide, left, top, width, height, text, fill_color, text_color=TEXT):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = text_color


def add_process_flow(slide, items, left=Inches(0.6), top=Inches(2.15), box_w=Inches(2.2), box_h=Inches(0.95), gap=Inches(0.22)):
    x = left
    centers = []
    for idx, item in enumerate(items):
        fill = HILITE_BG if idx < len(items) - 1 else WARN_BG
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, top, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = BOX_LINE
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = item
        r.font.name = FONT
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = TEXT
        centers.append((x + box_w / 2, top + box_h / 2))
        x += box_w + gap

    for i in range(len(centers) - 1):
        x1 = centers[i][0] + box_w / 2 - Inches(0.05)
        x2 = centers[i + 1][0] - box_w / 2 + Inches(0.05)
        y = centers[i][1]
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y, x2, y)
        conn.line.color.rgb = ACCENT
        conn.line.width = Pt(2.25)


def add_comparison_table(slide, left, top):
    headers = ["路径", "现有工作主问题", "我们主问题"]
    rows = [
        ["representative sampling", "模哪些样本 / 模多少样本", "workload 如何形成 simulator 分析对象"],
        ["proxy benchmark / motif", "用什么小程序代表原 workload", "原 workload 内部如何形成结构化归属"],
        ["single-kernel profiling", "这个 kernel 为什么慢", "整条 workload 如何组织成 reasoning units"],
    ]
    widths = [Inches(2.6), Inches(3.75), Inches(4.45)]
    row_h = Inches(0.7)

    x = left
    y = top
    for idx, head in enumerate(headers):
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, widths[idx], row_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = TITLE_BG
        rect.line.color.rgb = BG
        add_text_box(slide, x + Inches(0.08), y + Inches(0.14), widths[idx] - Inches(0.16), Inches(0.3),
                     head, font_size=15, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
        x += widths[idx]

    for row_idx, row in enumerate(rows):
        x = left
        y = top + row_h + row_idx * Inches(0.82)
        for col_idx, cell in enumerate(row):
            rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, widths[col_idx], Inches(0.82))
            rect.fill.solid()
            rect.fill.fore_color.rgb = BOX_BG if col_idx < 2 else HILITE_BG
            rect.line.color.rgb = BOX_LINE
            add_text_box(slide, x + Inches(0.08), y + Inches(0.12), widths[col_idx] - Inches(0.16), Inches(0.56),
                         cell, font_size=14, bold=(col_idx == 0), align=PP_ALIGN.CENTER)
            x += widths[col_idx]


def add_primitive_boxes(slide, primitives, left, top, width, height, cols=3):
    gap_x = Inches(0.18)
    gap_y = Inches(0.18)
    box_w = (width - gap_x * (cols - 1)) / cols
    rows = (len(primitives) + cols - 1) // cols
    box_h = (height - gap_y * (rows - 1)) / rows
    for idx, primitive in enumerate(primitives):
        row = idx // cols
        col = idx % cols
        x = left + col * (box_w + gap_x)
        y = top + row * (box_h + gap_y)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = HILITE_BG
        shape.line.color.rgb = BOX_LINE
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = primitive
        r.font.name = FONT
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = TEXT


def add_two_case_cards(slide):
    cases = [
        (
            "Case A: gemm_tiled vs attention_score",
            [
                "说明不能只按算子名字拆分",
                "二者存在共享执行模板，但语义不同",
                "问题被逼到 execution primitive 层",
            ],
            HILITE_BG,
        ),
        (
            "Case B: softmax_kernel vs context_mul",
            [
                "说明不能只按上层 attention 模块合并",
                "softmax 更像 reduction / normalize",
                "context_mul 更像 weighted aggregation",
            ],
            WARN_BG,
        ),
    ]
    lefts = [Inches(0.75), Inches(6.8)]
    for idx, (title, bullets, fill) in enumerate(cases):
        add_bullet_panel(slide, lefts[idx], Inches(1.75), Inches(5.75), Inches(3.35), title, bullets, fill_color=fill)


def add_three_value_cards(slide):
    titles = [
        ("定义价值", ["把经验性的 workload 映射过程显式化", "让 simulator 前的结构层变得可讨论"]),
        ("方法价值", ["把分析对象提升到 primitive / family 层", "不再停留在单 kernel 局部观察"]),
        ("扩展价值", ["为 family-aware simulator tuning 提供组织基础", "为后续 validation lane 留出稳定接口"]),
    ]
    left = Inches(0.75)
    width = Inches(3.95)
    gap = Inches(0.25)
    for idx, (title, bullets) in enumerate(titles):
        add_bullet_panel(slide, left + idx * (width + gap), Inches(2.05), width, Inches(2.65), title, bullets)


def add_next_steps(slide):
    steps = [
        ("Step 1", "做硬 transformer 主链 primitive 判据"),
        ("Step 2", "形成 family selection / boundary protocol"),
        ("Step 3", "再向 simulator lane 与 tuning lane 对接"),
    ]
    top = Inches(2.0)
    left = Inches(0.95)
    width = Inches(3.8)
    gap = Inches(0.4)
    for idx, (step, text) in enumerate(steps):
        x = left + idx * (width + gap)
        add_banner(slide, x, top, width, Inches(0.5), step, ACCENT if idx == 0 else ACCENT_2, RGBColor(255, 255, 255))
        add_bullet_panel(slide, x, top + Inches(0.65), width, Inches(1.9), text, [], fill_color=BOX_BG)
        inner = slide.shapes.add_textbox(x + Inches(0.15), top + Inches(1.0), width - Inches(0.3), Inches(1.1))
        tf = inner.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = TEXT
        if idx < 2:
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                x + width + Inches(0.08),
                top + Inches(1.55),
                x + width + gap - Inches(0.08),
                top + Inches(1.55),
            )
            conn.line.color.rgb = ACCENT
            conn.line.width = Pt(2.25)


def add_summary_points(slide):
    bullets = [
        "当前空缺真实存在：workload 到 simulator 之间缺少结构化接口",
        "我们的方法方向明确：execution primitive -> family",
        "Transformer 主链已经提供了第一版可落地原型",
    ]
    add_bullet_panel(slide, Inches(0.9), Inches(1.9), Inches(7.1), Inches(3.25), "三点总结", bullets)
    add_banner(
        slide,
        Inches(8.35),
        Inches(2.2),
        Inches(4.0),
        Inches(1.9),
        "核心结论\n不是减少几个 simulation samples，\n而是定义一层新的结构化中间表示。",
        HILITE_BG,
        TEXT,
    )


def build_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(2.0))
    hero.fill.solid()
    hero.fill.fore_color.rgb = TITLE_BG
    hero.line.color.rgb = TITLE_BG
    add_text_box(slide, Inches(0.7), Inches(0.62), Inches(10.8), Inches(0.6),
                 "从 Workload 到 Simulator 的结构化 Family 接口",
                 font_size=30, bold=True, color=RGBColor(255, 255, 255))
    add_text_box(slide, Inches(0.72), Inches(1.28), Inches(8.8), Inches(0.3),
                 "导师汇报初版 | 方法论文主线 | mini_transformer_v4 原型",
                 font_size=14, color=RGBColor(219, 227, 236))
    add_bullet_panel(
        slide, Inches(0.9), Inches(2.55), Inches(5.65), Inches(2.45), "本次汇报只做两件事",
        [
            "证明当前 workload -> simulator reasoning 之间存在真实结构空缺",
            "证明 execution primitive -> family 不是空想，而是一个可落地的原型方向",
        ]
    )
    add_banner(slide, Inches(7.1), Inches(2.75), Inches(5.0), Inches(1.75),
               "主线定位\n方法论文，而不是采样技巧报告", WARN_BG)
    add_footer(slide, "2026-04-19 | modern-gpu-simulator-micro-2025")

    # Slide 2
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_top_bar(slide, "问题：现有流程仍缺少结构化接口", "核心不是 workload 不能模拟，而是 workload 很难被自然组织成 simulator 可承接的对象")
    add_bullet_panel(
        slide, Inches(0.75), Inches(1.55), Inches(4.15), Inches(3.8), "当前痛点",
        [
            "kernel 很多、phase 很杂、解释对象不显式",
            "同一 workload 往往依赖人工经验拆热点、并 case、选验证主线",
            "模拟器能跑 workload，但前面的组织层仍然模糊",
        ]
    )
    add_bullet_panel(
        slide, Inches(5.15), Inches(1.55), Inches(3.15), Inches(3.8), "现象",
        [
            "sample 可以减少",
            "proxy 可以构造",
            "single kernel 可以分析",
            "但中间结构层仍然缺位",
        ],
        fill_color=HILITE_BG,
    )
    add_banner(slide, Inches(8.55), Inches(1.85), Inches(4.05), Inches(2.9),
               "结论\n现有流程能模拟 workload，\n但还不能自然地组织 workload。", WARN_BG)
    add_footer(slide, "第 2 页 | 问题定义")

    # Slide 3
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_top_bar(slide, "相关工作一：PKA 做了什么", "Principal Kernel Analysis 代表的是 sampled GPU simulation 里的 representative-kernel 路线")
    add_image_card(
        slide,
        Inches(0.7),
        Inches(1.45),
        Inches(12.0),
        Inches(4.75),
        ASSET_DIR / "pka_method_crop.png",
        "PKA：通过选择代表 kernel 和 thread block 来减少模拟量",
        [
            "核心目标：减少需要进入 simulator 的 kernel 和 thread blocks",
            "核心方法：PCA + 聚类选 representative kernels，再做 kernel projection",
            "本质上是在压缩 simulation samples，而不是显式定义 workload 的结构接口",
        ],
        "图源：Baddouh et al., MICRO 2021, Principal Kernel Analysis",
    )
    add_footer(slide, "第 3 页 | related work 1")

    # Slide 4
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_top_bar(slide, "相关工作二：Sieve 做了什么", "Sieve 代表的是更稳健的 stratified sampling 路线，重点仍然是 sampled simulation")
    add_image_card(
        slide,
        Inches(0.7),
        Inches(1.45),
        Inches(12.0),
        Inches(4.75),
        ASSET_DIR / "sieve_method_crop.png",
        "Sieve：先分层，再选 representative invocations 做加权预测",
        [
            "核心目标：提高 sampled simulation 的稳定性与精度",
            "核心方法：profile -> stratification -> representative kernel set -> weighted prediction",
            "改进了 sampling 质量，但目标仍然是 workload estimation / sampled simulation",
        ],
        "图源：Nalbantov et al., ISPASS 2023, Sieve",
    )
    add_footer(slide, "第 4 页 | related work 2")

    # Slide 5
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_top_bar(slide, "我们的区别：不是采样技巧，而是结构化接口", "我们不是回答“模哪些样本”，而是回答“按什么结构去解释和组织 workload”")
    add_comparison_table(slide, Inches(0.75), Inches(1.55))
    add_banner(slide, Inches(0.9), Inches(5.0), Inches(11.5), Inches(0.78),
               "PKA / Sieve 压缩的是 simulation samples；我们试图补上的是 workload -> simulator reasoning 的结构层。", WARN_BG)
    add_footer(slide, "第 5 页 | 我们和 related work 的区别")

    # Slide 6
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_top_bar(slide, "方法价值与可行性", "这一页回答两个问题：为什么值得做，为什么不是异想天开")
    add_bullet_panel(
        slide, Inches(0.75), Inches(1.55), Inches(3.8), Inches(3.6), "价值",
        [
            "把 workload 到 simulator 的经验性映射过程显式化",
            "把分析对象提升到 primitive / family 层，而不是只盯单个 kernel",
            "为后续 family-aware simulator tuning 提供组织基础",
        ]
    )
    add_bullet_panel(
        slide, Inches(4.8), Inches(1.55), Inches(3.8), Inches(3.6), "可行性",
        [
            "我们不是按算子名分组，也不是按上层模块分组",
            "判据是被 boundary case 反例逼出来的，不是拍脑袋",
            "Transformer 主链已经能覆盖多种 execution primitive",
        ],
        fill_color=HILITE_BG,
    )
    add_banner(slide, Inches(8.85), Inches(1.9), Inches(3.55), Inches(2.8),
               "核心判断\n这个方向值得做，\n因为当前空缺真实存在；\n这个方向也可行，\n因为原型已经能长出来。", WARN_BG)
    add_footer(slide, "第 6 页 | 价值与可行性")

    # Slide 7
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_top_bar(slide, "Transformer 主计算路线：我们现在具体做了什么", "第一版不是泛泛谈 family，而是先在一条完整主路径上压出 execution primitives 和 family 线索")
    add_process_flow(
        slide,
        [
            "workload trace",
            "execution primitive",
            "family",
            "simulator reasoning lane",
        ],
        left=Inches(0.7),
        top=Inches(1.55),
        box_w=Inches(2.8),
        box_h=Inches(1.0),
        gap=Inches(0.28),
    )
    add_process_flow(
        slide,
        [
            "QKV / projection",
            "attention_score",
            "softmax",
            "context_mul",
            "residual / norm / FFN",
        ],
        left=Inches(0.55),
        top=Inches(3.0),
        box_w=Inches(2.35),
        box_h=Inches(0.95),
        gap=Inches(0.15),
    )
    add_text_box(slide, Inches(0.78), Inches(4.35), Inches(2.7), Inches(0.3), "这条主链已覆盖的 execution primitives", font_size=17, bold=True)
    add_primitive_boxes(
        slide,
        [
            "Dense Tile Compute",
            "Pairwise Score",
            "Reduction / Normalize",
            "Weighted Aggregation",
            "Elementwise / Fusion",
        ],
        Inches(0.8),
        Inches(4.72),
        Inches(11.7),
        Inches(1.28),
        cols=3,
    )
    add_footer(slide, "第 7 页 | Transformer 主链原型")

    # Slide 8
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_top_bar(slide, "当前原型：Boundary Case 如何逼出判据", "这一步最关键的作用，是证明 family 的定义必须下沉到 execution primitive，而不是停留在算子语义")
    add_two_case_cards(slide)
    add_banner(slide, Inches(1.25), Inches(5.45), Inches(10.8), Inches(0.8),
               "结论：family 的核心不是语义标签，而是 execution primitive。", WARN_BG)
    add_footer(slide, "第 8 页 | 当前原型")

    # Slide 9
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_top_bar(slide, "下一步：先把方法做硬，再向 Simulator 对接", "当前阶段的重点不是铺更多 workload，而是先把方法论原型打硬")
    add_next_steps(slide)
    add_banner(slide, Inches(1.05), Inches(5.15), Inches(11.0), Inches(0.78),
               "当前优先级：Transformer 主链 primitive 判据 > family selection protocol > simulator lane 对接", HILITE_BG)
    add_footer(slide, "第 9 页 | 下一步")

    # Slide 10
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_top_bar(slide, "总结", "本次汇报的目标不是证明所有实验已经完成，而是先把论文问题、方法接口和第一版原型讲清楚")
    add_summary_points(slide)
    add_footer(slide, "第 10 页 | 总结")

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"saved to {OUT_PATH}")


if __name__ == "__main__":
    build_deck()
