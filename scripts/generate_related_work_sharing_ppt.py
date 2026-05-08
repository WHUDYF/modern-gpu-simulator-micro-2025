import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_OUTPUT = REPO_ROOT / "docs" / "slides" / "related_work_sharing_2026-04-25.pptx"


def parse_args(argv=None):
    parser = argparse.ArgumentParser(description="Generate the related-work sharing slide deck.")
    parser.add_argument(
        "--template",
        type=Path,
        default=None,
        help="Optional PowerPoint template path. If omitted, python-pptx creates a blank deck.",
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=DEFAULT_OUTPUT,
        help=f"Output .pptx path. Defaults to {DEFAULT_OUTPUT}.",
    )
    return parser.parse_args(argv)


def load_presentation(template: Path | None):
    if template is None:
        return Presentation()
    template = template.expanduser().resolve()
    if not template.exists():
        raise FileNotFoundError(f"template not found: {template}")
    return Presentation(str(template))


def set_run_style(run, size, bold=False, color=(0, 0, 0), font="Aptos"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = RGBColor(*color)


def add_footer(slide, text="Yifu Ding, Apr-25-2026"):
    box = slide.shapes.add_textbox(Inches(10.0), Inches(7.1), Inches(2.2), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    set_run_style(run, 12, color=(90, 90, 90))


def add_header_bar(slide, title):
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.55)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(22, 78, 99)
    bar.line.color.rgb = RGBColor(22, 78, 99)
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.1), Inches(11.8), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_run_style(run, 24, bold=True, color=(255, 255, 255))


def add_bullets(slide, items, x=0.85, y=1.2, w=11.2, h=5.8, font_size=24, level0_color=(0, 0, 0)):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = text
        color = level0_color if level == 0 else (70, 70, 70)
        set_run_style(run, font_size - level * 2, bold=(level == 0), color=color)
    return box


def add_two_col_callout(slide, left_title, left_items, right_title, right_items):
    left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(1.3), Inches(5.9), Inches(4.8))
    left.fill.solid()
    left.fill.fore_color.rgb = RGBColor(241, 248, 250)
    left.line.color.rgb = RGBColor(180, 205, 214)
    lt = slide.shapes.add_textbox(Inches(0.95), Inches(1.55), Inches(5.2), Inches(0.4))
    p = lt.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = left_title
    set_run_style(run, 22, bold=True, color=(22, 78, 99))
    add_bullets(slide, left_items, x=0.95, y=2.0, w=5.1, h=3.8, font_size=19)

    right = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.3), Inches(5.9), Inches(4.8))
    right.fill.solid()
    right.fill.fore_color.rgb = RGBColor(252, 247, 239)
    right.line.color.rgb = RGBColor(212, 188, 149)
    rt = slide.shapes.add_textbox(Inches(7.1), Inches(1.55), Inches(5.2), Inches(0.4))
    p = rt.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = right_title
    set_run_style(run, 22, bold=True, color=(140, 82, 35))
    add_bullets(slide, right_items, x=7.1, y=2.0, w=5.1, h=3.8, font_size=19)


def add_timeline_slide(slide):
    add_header_bar(slide, "这次汇报主要分享什么")
    add_bullets(
        slide,
        [
            "我的 idea 还没有完全做完，所以这次先不强调方法结果，而是分享这几篇对我影响最大的参考论文。",
            "我主要想回答三个问题：",
            ("1. sampled GPU simulation 这条线，最近到底在补哪些维度？", 1),
            ("2. 这些论文各自停在了哪里？", 1),
            ("3. 我们后续该吸收它们的哪一部分，而不是重复它们？", 1),
        ],
        x=0.9,
        y=1.1,
        w=11.3,
        h=1.9,
        font_size=24,
    )

    papers = [
        ("MICRO 2021", "PKA", "behavior feature space"),
        ("ISPASS 2023", "Sieve", "work-size stratification"),
        ("MICRO 2023", "Photon", "online execution-path structure"),
        ("MICRO 2025", "STEM+ROOT", "runtime heterogeneity"),
        ("arXiv 2026", "GCL-Sampler", "learned similarity"),
    ]
    x0 = 0.9
    y = 3.3
    width = 2.2
    gap = 0.18
    for idx, (year, name, dim) in enumerate(papers):
        x = Inches(x0 + idx * (width + gap))
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, Inches(y), Inches(width), Inches(2.05))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(236, 243, 247)
        shape.line.color.rgb = RGBColor(163, 191, 204)
        tb = slide.shapes.add_textbox(x + Inches(0.12), Inches(y + 0.18), Inches(width - 0.24), Inches(1.7))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = year
        set_run_style(r, 15, bold=True, color=(22, 78, 99))
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = name
        set_run_style(r, 20, bold=True, color=(0, 0, 0))
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = dim
        set_run_style(r, 13, color=(80, 80, 80))


def add_paper_slide(slide, title, venue, target, method, takeaway, color):
    add_header_bar(slide, title)
    venue_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.95), Inches(2.6), Inches(0.55))
    venue_box.fill.solid()
    venue_box.fill.fore_color.rgb = RGBColor(*color)
    venue_box.line.color.rgb = RGBColor(*color)
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.08), Inches(2.2), Inches(0.25))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = venue
    set_run_style(r, 16, bold=True, color=(255, 255, 255))

    add_two_col_callout(
        slide,
        "这篇工作解决什么",
        [target, "核心目标仍然是 sampled simulation / representative compression。"],
        "我最关心的方法点",
        method,
    )

    note = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.25), Inches(11.8), Inches(0.55))
    note.fill.solid()
    note.fill.fore_color.rgb = RGBColor(248, 250, 252)
    note.line.color.rgb = RGBColor(196, 210, 219)
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(6.35), Inches(11.2), Inches(0.25))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "对我们的直接启发： " + takeaway
    set_run_style(r, 18, bold=True, color=(45, 45, 45))


def main(argv=None):
    args = parse_args(argv)
    prs = load_presentation(args.template)
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    blank = prs.slide_layouts[0]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    title_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.05))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(22, 78, 99)
    title_bar.line.color.rgb = RGBColor(22, 78, 99)
    tb = slide.shapes.add_textbox(Inches(0.72), Inches(1.45), Inches(11.6), Inches(1.4))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "参考论文分享：从 workload compression 到 simulator 结构层"
    set_run_style(r, 30, bold=True, color=(20, 20, 20))
    p = tb.text_frame.add_paragraph()
    r = p.add_run()
    r.text = "这次汇报先不讲最终方法结果，重点讲几篇真正影响我当前思路的工作。"
    set_run_style(r, 21, color=(70, 70, 70))
    sub = slide.shapes.add_textbox(Inches(0.75), Inches(4.85), Inches(4.8), Inches(0.4))
    p = sub.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Yifu Ding, Apr-25-2026"
    set_run_style(r, 18, color=(80, 80, 80))
    tagline = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(3.45), Inches(8.8), Inches(0.7))
    tagline.fill.solid()
    tagline.fill.fore_color.rgb = RGBColor(241, 248, 250)
    tagline.line.color.rgb = RGBColor(180, 205, 214)
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(3.63), Inches(8.2), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "主线问题：现有工作擅长“压缩谁、采样谁”，但很少继续回答“压缩之后如何组织成 simulator 分析对象”。"
    set_run_style(r, 18, color=(30, 30, 30))

    # Slide 2
    slide = prs.slides.add_slide(blank)
    add_header_bar(slide, "目录")
    add_bullets(
        slide,
        [
            "1. 这次为什么先汇报论文，而不是汇报最终 idea",
            "2. 五篇我重点参考的工作：PKA / Sieve / Photon / STEM+ROOT / GCL-Sampler",
            "3. 每篇工作到底补了哪一个维度",
            "4. 我现在从这些论文里真正想吸收什么",
            "5. 对我们后续 PPT 和方法线的影响",
        ],
        x=1.1,
        y=1.45,
        w=10.8,
        h=4.6,
        font_size=24,
    )
    add_footer(slide)

    # Slide 3
    slide = prs.slides.add_slide(blank)
    add_timeline_slide(slide)
    add_footer(slide)

    # Slide 4
    slide = prs.slides.add_slide(blank)
    add_paper_slide(
        slide,
        "PKA：behavior feature space 作为 representative compression 的基础",
        "MICRO 2021",
        "从大量 kernel 中选出少量 representative kernels，以降低大规模 GPU workload 的 simulation cost。",
        [
            "采集 microarchitecture-independent behavior features。",
            "通过 PCA + K-means 在行为特征空间中组织对象。",
            "最后输出 representative kernels + clusters。",
        ],
        "A 线的主分组结构应该尽量对齐 PKA 的行为特征空间；PKA 更适合作为 frontend anchor，而不是我们去正面重做的对象。",
        (37, 99, 133),
    )
    add_footer(slide)

    # Slide 5
    slide = prs.slides.add_slide(blank)
    add_paper_slide(
        slide,
        "Sieve：work-size 不是附属信息，而是 grouping 的必要条件",
        "ISPASS 2023",
        "通过 stratified sampling 让 strata 内 execution time variance 更小，从而提高 sampled simulation 的稳定性。",
        [
            "把 instruction count 当作 work-size proxy。",
            "不仅看行为相似，还显式控制工作量尺度。",
            "代表对象不只要“像”，还要在规模上接近。",
        ],
        "A 线后续不能只给出“谁代表谁”，还要带出 member_count、time_weight、workload_scale；否则 B 线的重要性判断没有基础。",
        (91, 110, 64),
    )
    add_footer(slide)

    # Slide 6
    slide = prs.slides.add_slide(blank)
    add_paper_slide(
        slide,
        "Photon：在线执行路径结构本身也可以成为特征",
        "MICRO 2023",
        "在不依赖重 upfront profiling 的情况下，在线决定 sampled simulation 应采用哪一层粒度。",
        [
            "在线分析 kernel / warp / basic-block。",
            "构造 GPU BBV，按执行路径结构自适应切换采样粒度。",
            "提醒我们：family / regime 的证据源不一定只有离线 counters。",
        ],
        "Photon 暂时不是我们第一版必须实现的部分，但它说明后续若要增强 family 判据，可以考虑引入在线执行路径结构。",
        (134, 82, 32),
    )
    add_footer(slide)

    # Slide 7
    slide = prs.slides.add_slide(blank)
    add_paper_slide(
        slide,
        "STEM+ROOT：同名 kernel 也可能高度 heterogeneous",
        "MICRO 2025",
        "解决 invocation 级 runtime distribution heterogeneity，让 sampled simulation 更准、更可控。",
        [
            "ROOT 递归拆分 heterogeneous clusters。",
            "STEM 根据 CoV 和 error bound 分配 sample budget。",
            "核心提醒：grouping 不能只看结构相似，还要看 runtime distribution。",
        ],
        "这篇工作最直接影响我对 A 线和 B 线的理解：前端代表对象必须带出时间分布和权重信息，否则后面会把少量异质对象错当主干对象。",
        (118, 62, 138),
    )
    add_footer(slide)

    # Slide 8
    slide = prs.slides.add_slide(blank)
    add_header_bar(slide, "把这几篇论文放在一起看，我目前的统一理解")
    add_bullets(
        slide,
        [
            "PKA 告诉我：必须先在 behavior feature space 里做 representative compression。",
            "Sieve 告诉我：只看行为相似还不够，work-size 也必须进来。",
            "Photon 告诉我：在线执行路径结构可以成为后续更强的证据源。",
            "STEM+ROOT 告诉我：同名或同结构对象也可能在 runtime distribution 上高度异质。",
            "GCL-Sampler 则提醒我：如果只是继续卷“更强的相似性发现”，竞争会很激烈。",
        ],
        x=0.95,
        y=1.2,
        w=11.2,
        h=3.5,
        font_size=22,
    )
    add_two_col_callout(
        slide,
        "所以我现在不想重复做的事",
        [
            "继续把工作写成“更强的 representative kernel selection”。",
            "继续把方法停在 sampled simulation 的 sample selection 层。",
        ],
        "所以我现在想保住的边界",
        [
            "把前端压缩后的对象继续提升成 simulator-side structured objects。",
            "重点放在 compression 之后、simulator 之前的结构层。",
        ],
    )
    add_footer(slide)

    # Slide 9
    slide = prs.slides.add_slide(blank)
    add_header_bar(slide, "这些论文现在如何影响我们的修改方向")
    add_bullets(
        slide,
        [
            "A 线：主分组结构应尽量对齐 PKA 的 behavior feature space。",
            "A 线：还要吸收 Sieve / STEM 的权重表达，保留 member_count、time_weight、workload_scale。",
            "B 线：family 主要承接硬件分组，regime 再把算法功能角色和上下文引入进来。",
            "B 线：不要退化成“按算子名贴标签”，而要区分 hardware grouping 和 algorithm-function grouping。",
            "整体边界：我们不是继续卷 sample selection，而是构建从 compressed kernels 到 simulator reasoning objects 的接口层。",
        ],
        x=0.95,
        y=1.25,
        w=11.2,
        h=4.8,
        font_size=21,
    )
    add_footer(slide)

    # Slide 10
    slide = prs.slides.add_slide(blank)
    add_header_bar(slide, "我希望这次汇报带走的结论")
    add_bullets(
        slide,
        [
            "1. sampled GPU simulation 这条线并没有停，2021-2026 一直在快速推进。",
            "2. 这些工作分别补上了 behavior、work-size、path structure、runtime heterogeneity、learned similarity 五类维度。",
            "3. 我们当前最值得做的，不是再造一个更强 sampling 方法，而是在这些前端压缩结果之后补上 simulator 结构层。",
            "4. 因为 idea 还没完全做完，这次我更想把参考论文和吸收路径讲清楚，而不是过早宣称最终方法已经定型。",
        ],
        x=0.9,
        y=1.35,
        w=11.3,
        h=4.5,
        font_size=22,
    )
    add_footer(slide)

    output = args.output.expanduser().resolve()
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))


if __name__ == "__main__":
    main()
